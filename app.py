from flask import Flask, render_template, request, jsonify, session, redirect, url_for, send_from_directory, flash
import os
import random
from functools import wraps
import time
from werkzeug.utils import secure_filename
import json
import PyPDF2
from datetime import datetime
import io
from pptx import Presentation
from PIL import Image

app = Flask(__name__)
app.secret_key = os.urandom(24)

# 數據文件路徑
CURRENT_DIR = os.path.dirname(os.path.abspath(__file__))
PRESENTATIONS_FILE = os.path.join(CURRENT_DIR, 'static', 'data', 'presentations.json')
UPLOADS_FILE = os.path.join(CURRENT_DIR, 'static', 'data', 'uploads.json')
QUESTIONS_FILE = os.path.join(CURRENT_DIR, 'static', 'data', 'questions.json')
REFERENCES_FILE = os.path.join(CURRENT_DIR, 'static', 'data', 'references.json')
SUPPLEMENTARY_FILE = os.path.join(CURRENT_DIR, 'static', 'data', 'supplementary.json')
SLIDES_DIR = os.path.join(CURRENT_DIR, 'static', 'slides')
UPLOAD_FOLDER = os.path.join(CURRENT_DIR, 'static', 'uploads')

# 確保目錄存在
os.makedirs(os.path.join(CURRENT_DIR, 'static', 'data'), exist_ok=True)
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(SLIDES_DIR, exist_ok=True)

# 從文件加載數據
def load_data_from_file(file_path, default_data=None):
    if default_data is None:
        default_data = []
    
    if os.path.exists(file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            print(f"載入數據時發生錯誤: {str(e)}")
            return default_data
    else:
        return default_data

# 保存數據到文件
def save_data_to_file(file_path, data):
    try:
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=4)
        return True
    except Exception as e:
        print(f"保存數據時發生錯誤: {str(e)}")
        return False

# 將 PowerPoint 轉換為圖片
def convert_ppt_to_images(ppt_path, presentation_id):
    """
    將 PowerPoint 文件轉換為圖片並保存
    使用 python-pptx 和 PIL 進行轉換，提供更好的渲染效果
    """
    try:
        print(f"開始轉換 PowerPoint 文件: {ppt_path}")
        
        # 創建保存圖片的目錄
        slides_dir = os.path.join(app.static_folder, 'slides', f"presentation-{presentation_id}")
        os.makedirs(slides_dir, exist_ok=True)
        
        print(f"圖片將保存到: {slides_dir}")
        
        # 打開 PowerPoint 文件
        prs = Presentation(ppt_path)
        
        # 獲取頁數
        num_slides = len(prs.slides)
        print(f"PowerPoint 共有 {num_slides} 頁")
        
        # 設置圖片尺寸 (16:9 比例)
        width = 1280
        height = 720
        
        # 轉換每一頁為圖片
        for i, slide in enumerate(prs.slides):
            # 保存路徑
            img_path = os.path.join(slides_dir, f"slide-{i+1}.png")
            
            # 創建一個空白圖片
            img = Image.new('RGB', (width, height), 'white')
            
            # 使用 PIL 繪製文字
            from PIL import ImageDraw, ImageFont
            draw = ImageDraw.Draw(img)
            
            try:
                # 嘗試加載字體，如果失敗則使用默認字體
                try:
                    # 嘗試使用支持中文的字體
                    fonts_to_try = [
                        "simhei.ttf",  # 中文黑體
                        "simsun.ttc",  # 中文宋體
                        "msyh.ttc",    # 微軟雅黑
                        "arial.ttf",   # 英文字體
                    ]
                    
                    title_font = None
                    for font_name in fonts_to_try:
                        try:
                            title_font = ImageFont.truetype(font_name, 36)
                            break
                        except:
                            continue
                    
                    if title_font is None:
                        title_font = ImageFont.load_default()
                        
                    content_font = ImageFont.truetype(font_name, 24) if title_font != ImageFont.load_default() else title_font
                    small_font = ImageFont.truetype(font_name, 18) if title_font != ImageFont.load_default() else title_font
                except:
                    title_font = ImageFont.load_default()
                    content_font = title_font
                    small_font = title_font
                
                # 繪製幻燈片背景
                draw.rectangle([(0, 0), (width, 80)], fill="#4472C4")  # 頂部藍色條
                
                # 提取和繪製幻燈片內容
                y_position = 100
                title_found = False
                
                # 處理每個形狀
                for shape in slide.shapes:
                    if not hasattr(shape, 'text'):
                        continue
                        
                    text = shape.text.strip()
                    if not text:
                        continue
                    
                    # 檢查形狀類型和位置來確定它是標題還是內容
                    if not title_found and (hasattr(shape, 'is_title') and shape.is_title) or y_position < 150:
                        # 這是標題
                        draw.text((40, 20), text, fill="white", font=title_font)
                        title_found = True
                    else:
                        # 這是內容
                        # 處理長文本，自動換行
                        words = text.split()
                        lines = []
                        current_line = ""
                        
                        for word in words:
                            test_line = current_line + " " + word if current_line else word
                            # 檢查行寬度
                            if draw.textlength(test_line, font=content_font) < width - 80:
                                current_line = test_line
                            else:
                                lines.append(current_line)
                                current_line = word
                        
                        if current_line:
                            lines.append(current_line)
                        
                        # 繪製文本行
                        for line in lines:
                            draw.text((40, y_position), line, fill="black", font=content_font)
                            y_position += 30
                        
                        y_position += 20  # 段落間距
                
                # 如果沒有找到標題，添加默認標題
                if not title_found:
                    draw.text((40, 20), f"幻燈片 {i+1}", fill="white", font=title_font)
                
                # 繪製頁碼在右下角
                page_text = f"{i+1}/{num_slides}"
                draw.text((width-100, height-30), page_text, fill="black", font=small_font)
                
                # 嘗試提取和繪製圖片
                for shape in slide.shapes:
                    if hasattr(shape, 'image'):
                        try:
                            # 提取圖片數據
                            image_stream = io.BytesIO(shape.image.blob)
                            shape_img = Image.open(image_stream)
                            
                            # 計算圖片位置和大小
                            img_x = width // 2 - shape_img.width // 2
                            img_y = y_position
                            
                            # 確保圖片不會太大
                            max_img_width = width - 80
                            max_img_height = height - y_position - 50
                            
                            if shape_img.width > max_img_width or shape_img.height > max_img_height:
                                # 等比例縮放
                                ratio = min(max_img_width / shape_img.width, max_img_height / shape_img.height)
                                new_width = int(shape_img.width * ratio)
                                new_height = int(shape_img.height * ratio)
                                shape_img = shape_img.resize((new_width, new_height), Image.LANCZOS)
                                
                                img_x = width // 2 - new_width // 2
                            
                            # 將圖片粘貼到幻燈片上
                            img.paste(shape_img, (img_x, img_y))
                            y_position += shape_img.height + 20
                        except Exception as img_error:
                            print(f"處理幻燈片 {i+1} 中的圖片時發生錯誤: {str(img_error)}")
                
            except Exception as inner_e:
                print(f"繪製幻燈片 {i+1} 時發生錯誤: {str(inner_e)}")
                # 繪製錯誤信息
                draw.text((width//2 - 200, height//2), f"無法渲染幻燈片 {i+1}", fill="red", font=title_font)
                draw.text((width//2 - 200, height//2 + 50), str(inner_e), fill="red", font=small_font)
            
            # 保存圖片
            img.save(img_path, 'PNG')
            print(f"已保存第 {i+1} 頁到 {img_path}")
        
        return num_slides
    except Exception as e:
        print(f"轉換 PowerPoint 時發生錯誤: {str(e)}")
        return 0

# 默認數據
default_presentations = [
    {
        "id": 1,
        "title": "課程介紹",
        "description": "這是課程的介紹簡報",
        "file_type": "pptx",
        "path": "slides/presentation-1",
        "original_file": "uploads/intro.pptx",
        "total_pages": 5
    },
    {
        "id": 2,
        "title": "第一章：基礎概念",
        "description": "介紹基礎概念和術語",
        "file_type": "pdf",
        "path": "uploads/chapter1.pdf",
        "original_file": "uploads/chapter1.pdf",
        "total_pages": 10
    },
    {
        "id": 3,
        "title": "第二章：進階技巧",
        "description": "探討更複雜的技術和方法",
        "file_type": "pptx",
        "path": "slides/presentation-3",
        "original_file": "uploads/advanced.pptx",
        "total_pages": 8
    }
]

default_uploads = [
    {
        "id": 1,
        "title": "課程大綱",
        "description": "本課程的詳細大綱和學習目標",
        "file_type": "pdf",
        "path": "uploads/syllabus.pdf"
    },
    {
        "id": 2,
        "title": "實作範例",
        "description": "課程中提到的實作範例程式碼",
        "file_type": "zip",
        "path": "uploads/examples.zip"
    },
    {
        "id": 3,
        "title": "課程教學影片",
        "description": "課程的輔助教學影片",
        "file_type": "video",
        "path": "https://www.youtube.com/embed/dQw4w9WgXcQ"
    }
]

default_questions = [
    {
        "id": 1,
        "name": "王同學",
        "title": "如何下載課程資料？",
        "content": "我在尋找課程的補充資料，但不知道從哪裡可以下載到，請問應該怎麼操作？",
        "date": "2023/04/15 10:30",
        "answer": "您可以在「資源」頁面找到所有可下載的課程資料。點擊頁面上方的「資源」選項卡，然後在列表中找到您需要的資料進行下載。",
        "answer_date": "2023/04/15 14:45",
        "helpful_count": 5,
        "related_to": 1,
        "comments": [
            {
                "id": 1,
                "name": "李同學",
                "content": "謝謝講師的回答，我找到了！",
                "date": "2023/04/16 09:15",
                "is_lecturer": False
            }
        ]
    },
    {
        "id": 2,
        "name": "陳同學",
        "title": "考試範圍包含哪些內容？",
        "content": "想請問一下期中考試的範圍是什麼？需要特別準備哪些章節？",
        "date": "2023/04/17 14:45",
        "answer": "考試範圍包含課程中所有的章節，特別是第一章和第二章的內容。建議重點複習課程簡報中標記為「重要」的部分。",
        "answer_date": "2023/04/17 16:30",
        "helpful_count": 8,
        "related_to": 2,
        "comments": []
    },
    {
        "id": 3,
        "name": "林同學",
        "title": "是否有課後輔導時間？",
        "content": "我對某些概念還不太理解，想知道是否有額外的輔導時間可以請教老師？",
        "date": "2023/04/18 09:15",
        "answer": "是的，每週三下午 2-4 點有課後輔導時間，地點在 A307 教室。您也可以透過平台的問答功能隨時提問。",
        "answer_date": "2023/04/18 10:20",
        "helpful_count": 3,
        "related_to": None,
        "comments": []
    },
    {
        "id": 4,
        "name": "張同學",
        "title": "如何使用AI演示功能？",
        "content": "我看到平台有AI演示功能，但不太清楚如何操作，可以提供一些指引嗎？",
        "date": "2023/04/19 11:30",
        "answer": None,
        "answer_date": None,
        "helpful_count": 1,
        "related_to": 3,
        "comments": []
    }
]

default_references = [
    {
        "id": 1,
        "title": "程式設計導論",
        "url": "https://example.com/intro-programming",
        "date": "2023-01-01"
    },
    {
        "id": 2,
        "title": "資料結構與演算法",
        "url": "https://example.com/data-structures",
        "date": "2023-01-02"
    },
    {
        "id": 3,
        "title": "網頁開發基礎",
        "url": "https://example.com/web-development",
        "date": "2023-01-03"
    }
]

default_supplementary = [
    {
        "id": 1,
        "title": "程式設計練習題",
        "description": "課程相關的程式設計練習題和解答",
        "filename": "programming_exercises.pdf",
        "path": "uploads/supplementary_1_programming_exercises.pdf",
        "date": "2023-01-01"
    },
    {
        "id": 2,
        "title": "期中考試複習資料",
        "description": "期中考試的複習重點和範例題目",
        "filename": "midterm_review.pdf",
        "path": "uploads/supplementary_2_midterm_review.pdf",
        "date": "2023-01-02"
    },
    {
        "id": 3,
        "title": "專題製作指南",
        "description": "期末專題製作的詳細指南和評分標準",
        "filename": "project_guide.pdf",
        "path": "uploads/supplementary_3_project_guide.pdf",
        "date": "2023-01-03"
    }
]

# 加載數據
presentations = load_data_from_file(PRESENTATIONS_FILE, default_presentations)
uploads = load_data_from_file(UPLOADS_FILE, default_uploads)
questions = load_data_from_file(QUESTIONS_FILE, default_questions)
references = load_data_from_file(REFERENCES_FILE, default_references)
supplementary = load_data_from_file(SUPPLEMENTARY_FILE, default_supplementary)

# 允許的檔案類型
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif'}

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# 模擬講師帳號
lecturer_accounts = {
    "admin": {
        "password": "admin123",
        "name": "梁坤棠",
        "title": "副總經理"
    }
}

# 登入裝飾器
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'lecturer_username' not in session:
            return redirect(url_for('index'))
        return f(*args, **kwargs)
    return decorated_function

# 路由
@app.route('/')
def index():
    # 嘗試載入講師資料（無論是否登入）
    config_path = os.path.join(os.path.dirname(__file__), 'static', 'config', 'lecturer_info.json')
    
    if os.path.exists(config_path):
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                config = json.load(f)
                
                # 確保session中有講師資料
                if 'lecturer_name' not in session and 'lecturer_name' in config:
                    session['lecturer_name'] = config.get('lecturer_name')
                if 'lecturer_title' not in session and 'lecturer_title' in config:
                    session['lecturer_title'] = config.get('lecturer_title')
                if 'lecturer_email' not in session and 'lecturer_email' in config:
                    session['lecturer_email'] = config.get('lecturer_email')
        except Exception as e:
            print(f"載入講師資料時發生錯誤: {str(e)}")
    
    return render_template('index.html', presentations=presentations, questions=questions)

@app.route('/presentation')
def presentation():
    presentations = load_data_from_file(PRESENTATIONS_FILE, default_presentations)
    
    # 獲取請求的簡報ID
    presentation_id = request.args.get('id')
    page = request.args.get('page', 1, type=int)
    
    # 如果沒有指定ID，顯示第一個簡報
    current_presentation = None
    if presentation_id:
        # 查找對應ID的簡報
        for p in presentations:
            if str(p['id']) == str(presentation_id):
                current_presentation = p
                break
    elif presentations:
        current_presentation = presentations[0]
        presentation_id = current_presentation['id']
    
    # 如果找不到簡報，顯示空頁面
    if not current_presentation:
        return render_template('presentation.html', 
                               presentations=presentations, 
                               current_presentation=None,
                               current_page=1,
                               total_pages=0)
    
    # 確保頁碼在有效範圍內
    total_pages = current_presentation.get('total_pages', 1)
    if page < 1:
        page = 1
    elif page > total_pages:
        page = total_pages
    
    # 檢查是否為 PowerPoint 文件，如果是，檢查是否已轉換為圖片
    if current_presentation.get('file_type') == 'pptx':
        # 檢查轉換後的圖片目錄是否存在
        slides_dir = os.path.join(app.static_folder, 'slides', f"presentation-{current_presentation['id']}")
        if not os.path.exists(slides_dir) or len(os.listdir(slides_dir)) == 0:
            # 如果不存在，執行轉換
            ppt_path = os.path.join(app.static_folder, 'uploads', current_presentation['original_file'])
            if os.path.exists(ppt_path):
                num_slides = convert_ppt_to_images(ppt_path, current_presentation['id'])
                # 更新簡報頁數
                if num_slides > 0:
                    current_presentation['total_pages'] = num_slides
                    # 保存更新後的數據
                    save_data_to_file(PRESENTATIONS_FILE, presentations)
                    total_pages = num_slides
    
    print(f"渲染簡報頁面: ID={presentation_id}, 頁碼={page}, 總頁數={total_pages}")
    return render_template('presentation.html', 
                           presentations=presentations, 
                           current_presentation=current_presentation,
                           current_page=page,
                           total_pages=total_pages)

@app.route('/presentation/<int:presentation_id>', methods=['GET'])
def presentation_with_id(presentation_id):
    # 重定向到查詢參數格式的URL
    page = request.args.get('page', 1, type=int)
    return redirect(url_for('presentation', id=presentation_id, page=page))

@app.route('/ai-demo')
def ai_demo():
    return render_template('ai-demo.html', presentations=presentations)

@app.route('/resources')
def resources():
    return render_template('resources.html', presentations=presentations, uploads=uploads, references=references, supplementary=supplementary)

@app.route('/qa')
def qa():
    # 加載最新的數據
    global questions, presentations
    questions = load_data_from_file(QUESTIONS_FILE, default_questions)
    presentations = load_data_from_file(PRESENTATIONS_FILE, default_presentations)
    
    return render_template('qa.html', questions=questions, presentations=presentations)

# 講師登入
@app.route('/lecturer/login', methods=['POST'])
def lecturer_login():
    if request.is_json:
        data = request.get_json()
        username = data.get('username')
        password = data.get('password')
        remember_me = data.get('rememberMe', False)
    else:
        username = request.form.get('username')
        password = request.form.get('password')
        remember_me = request.form.get('rememberMe') == 'on'
    
    # 驗證帳號密碼
    if username in lecturer_accounts and lecturer_accounts[username]["password"] == password:
        # 登入成功，只設置username，不覆蓋姓名和職稱
        session['lecturer_username'] = username
        
        # 檢查session中是否已有從配置文件載入的資料
        # 如果沒有，才使用默認資料
        if 'lecturer_name' not in session:
            session['lecturer_name'] = lecturer_accounts[username]["name"]
        if 'lecturer_title' not in session:
            session['lecturer_title'] = lecturer_accounts[username]["title"]
        
        # 確保載入配置文件中的資料
        config_path = os.path.join(os.path.dirname(__file__), 'static', 'config', 'lecturer_info.json')
        
        if os.path.exists(config_path):
            try:
                with open(config_path, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    
                    # 使用配置文件中的資料
                    if 'lecturer_name' in config:
                        session['lecturer_name'] = config.get('lecturer_name')
                    if 'lecturer_title' in config:
                        session['lecturer_title'] = config.get('lecturer_title')
                    if 'lecturer_email' in config:
                        session['lecturer_email'] = config.get('lecturer_email')
                        print("登入時已從配置文件載入講師資料")
            except Exception as e:
                print(f"登入時載入講師資料錯誤: {str(e)}")
        
        if not remember_me:
            # 如果不記住我，設置session在瀏覽器關閉時過期
            session.permanent = False
        else:
            # 否則設置session持久化
            session.permanent = True
        
        return jsonify({"success": True})
    else:
        return jsonify({"success": False, "message": "帳號或密碼錯誤"})

# 講師登出
@app.route('/lecturer/logout')
def lecturer_logout():
    # 清除特定的session變數而不是整個session
    session.pop('lecturer_username', None)
    # 保留講師個人資料，這樣登出後重新登入時仍能保持個人化設定
    # session.pop('lecturer_name', None)
    # session.pop('lecturer_title', None)
    # session.pop('lecturer_email', None)
    return redirect(url_for('index'))

# 講師儀表板
@app.route('/lecturer/dashboard')
@login_required
def lecturer_dashboard():
    # 獲取最新的數據
    global questions, presentations, uploads, references, supplementary
    questions = load_data_from_file(QUESTIONS_FILE, default_questions)
    presentations = load_data_from_file(PRESENTATIONS_FILE, default_presentations)
    supplementary = load_data_from_file(SUPPLEMENTARY_FILE, default_supplementary)
    
    # 加載講師信息
    lecturer_info = lecturer_accounts.get(session.get('lecturer_username'))
    
    # 獲取未回答的問題數量
    unanswered_questions_count = sum(1 for q in questions if q.get('answer') is None)
    
    # 加載教材數據
    return render_template('lecturer/dashboard.html', 
                          lecturer=lecturer_info, 
                          presentations=presentations, 
                          uploads=uploads,
                          references=references,
                          supplementary=supplementary,
                          questions=questions,
                          unanswered_questions_count=unanswered_questions_count)

@app.route('/lecturer/answer-question', methods=['POST'])
@login_required
def lecturer_answer_question():
    try:
        data = request.get_json()
        question_id = data.get('question_id')
        answer = data.get('answer')
        
        # 驗證必要欄位
        if not question_id or not answer:
            return jsonify({"success": False, "message": "缺少必要欄位"})
        
        # 加載現有問題
        global questions
        questions = load_data_from_file(QUESTIONS_FILE, default_questions)
        
        # 查找問題並添加回答
        for question in questions:
            if question['id'] == int(question_id):
                question['answer'] = answer
                question['answer_date'] = datetime.now().strftime('%Y/%m/%d %H:%M')
                
                # 保存到文件
                success = save_data_to_file(QUESTIONS_FILE, questions)
                if not success:
                    return jsonify({"success": False, "message": "保存回答時發生錯誤"})
                
                return jsonify({"success": True})
        
        return jsonify({"success": False, "message": "找不到指定的問題"})
    except Exception as e:
        app.logger.error(f"回答問題時發生錯誤: {str(e)}")
        return jsonify({"success": False, "message": f"回答問題時發生錯誤: {str(e)}"})

# 講師個人資料更新
@app.route('/lecturer/update-profile', methods=['POST'])
def update_lecturer_profile():
    try:
        data = request.get_json()
        name = data.get('name', '')
        title = data.get('title', '')
        email = data.get('email', '')
        
        # 更新session中的講師資料
        session['lecturer_name'] = name
        session['lecturer_title'] = title
        session['lecturer_email'] = email
        
        # 將更新的資料保存到config.json文件中
        config_path = os.path.join(os.path.dirname(__file__), 'static', 'config', 'lecturer_info.json')
        
        # 確保目錄存在
        os.makedirs(os.path.dirname(config_path), exist_ok=True)
        
        # 讀取現有配置（如果存在）
        config = {}
        if os.path.exists(config_path):
            try:
                with open(config_path, 'r', encoding='utf-8') as f:
                    config = json.load(f)
            except:
                pass
        
        # 更新講師資料
        config['lecturer_name'] = name
        config['lecturer_title'] = title
        config['lecturer_email'] = email
        
        # 保存更新後的配置
        with open(config_path, 'w', encoding='utf-8') as f:
            json.dump(config, f, ensure_ascii=False, indent=4)
        
        return jsonify({'success': True, 'message': '個人資料已更新'})
    except Exception as e:
        print(f"更新個人資料時發生錯誤: {str(e)}")
        return jsonify({'success': False, 'message': f'更新個人資料時發生錯誤: {str(e)}'}), 500

# 講師照片更新
@app.route('/lecturer/update-photo', methods=['POST'])
def update_lecturer_photo():
    try:
        if 'photo' not in request.files:
            print("未提供照片文件")
            return jsonify({'success': False, 'message': '沒有提供照片'}), 400
        
        photo = request.files['photo']
        if photo.filename == '':
            print("未選擇檔案")
            return jsonify({'success': False, 'message': '未選擇檔案'}), 400
        
        if photo and allowed_file(photo.filename):
            # 使用固定檔案名稱
            filename = 'lecturer-photo.jpg'
            
            # 確保目錄存在
            upload_folder = os.path.join(os.path.dirname(__file__), 'static', 'images')
            os.makedirs(upload_folder, exist_ok=True)
            
            # 保存照片
            filepath = os.path.join(upload_folder, filename)
            print(f"正在保存照片到: {filepath}")
            photo.save(filepath)
            
            # 生成時間戳以防止瀏覽器緩存
            timestamp = int(time.time())
            session['photo_timestamp'] = timestamp
            
            # 返回照片URL (包含時間戳以防止緩存)
            photo_url = f"/static/images/{filename}?t={timestamp}"
            print(f"生成照片URL: {photo_url}")
            
            return jsonify({
                'success': True, 
                'message': '照片已更新',
                'photo_url': photo_url
            })
        else:
            print("不支援的檔案類型")
            return jsonify({'success': False, 'message': '不支援的檔案類型'}), 400
    except Exception as e:
        import traceback
        print(f"更新照片時發生錯誤: {str(e)}")
        print(traceback.format_exc())
        return jsonify({'success': False, 'message': f'更新照片時發生錯誤: {str(e)}'}), 500

# 上傳教材
@app.route('/lecturer/upload-material', methods=['POST'])
@login_required
def upload_material():
    if 'title' not in request.form:
        return jsonify({"success": False, "message": "未提供教材標題"})
    
    title = request.form.get('title')
    material_type = request.form.get('type')
    
    # 如果是文件上傳（非影片連結）
    if material_type != 'video':
        if 'file' not in request.files:
            return jsonify({"success": False, "message": "未提供檔案"})
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({"success": False, "message": "未選擇檔案"})
        
        # 檢查檔案類型
        if material_type == 'presentation' and not file.filename.lower().endswith(('.pdf', '.ppt', '.pptx')):
            return jsonify({"success": False, "message": "簡報必須是 PDF 或 PowerPoint 檔案 (.pdf, .ppt, .pptx)"})
        
        if material_type == 'document' and not file.filename.lower().endswith(('.pdf', '.doc', '.docx', '.txt')):
            return jsonify({"success": False, "message": "文件必須是 PDF、Word 或文字檔案 (.pdf, .doc, .docx, .txt)"})
        
        # 在真實環境中，這裡會保存文件到磁盤並更新資料庫
        try:
            # 安全的檔案名
            secure_filename = file.filename.replace(' ', '_')
            
            # 創建保存目錄（如果不存在）
            save_dir = os.path.join(os.path.dirname(__file__), 'static', 'uploads')
            os.makedirs(save_dir, exist_ok=True)
            
            # 保存文件
            file_path = os.path.join(save_dir, secure_filename)
            print(f"正在保存檔案到: {file_path}")
            file.save(file_path)
            
            # 獲取文件大小
            file_size = os.path.getsize(file_path)
            size_display = f"{file_size / 1024 / 1024:.1f} MB" if file_size > 1024 * 1024 else f"{file_size / 1024:.1f} KB"
            
            # 獲取當前日期
            from datetime import datetime
            current_date = datetime.now().strftime("%Y/%m/%d")
            
            # 獲取PDF頁數（如果是PDF文件）
            pages = 1
            if file.filename.lower().endswith('.pdf'):
                try:
                    with open(file_path, 'rb') as pdf_file:
                        pdf_reader = PyPDF2.PdfReader(pdf_file)
                        pages = len(pdf_reader.pages)
                except Exception as e:
                    print(f"無法讀取PDF頁數: {str(e)}")
                    pages = 20  # 如果無法讀取，使用默認值
            elif file.filename.lower().endswith(('.ppt', '.pptx')):
                # 對於PPT文件，目前無法直接讀取頁數，使用默認值
                pages = 20
                # 將PPT轉換為圖片
                pages = convert_ppt_to_images(file_path, len(presentations) + 1)
            
            # 將新教材添加到模擬數據
            if material_type == 'presentation':
                new_id = max([p["id"] for p in presentations]) + 1 if presentations else 1
                new_material = {
                    "id": new_id,
                    "title": title,
                    "filename": secure_filename,
                    "pages": pages,
                    "date": current_date,
                    "size": size_display,
                    "type": "presentation"
                }
                presentations.insert(0, new_material)
                save_data_to_file(PRESENTATIONS_FILE, presentations)
            else:
                # 添加到上傳檔案列表
                new_id = max([u["id"] for u in uploads]) + 1 if uploads else 1
                new_upload = {
                    "id": new_id,
                    "title": title,
                    "filename": secure_filename,
                    "type": material_type,
                    "date": current_date,
                    "size": size_display
                }
                uploads.insert(0, new_upload)
                save_data_to_file(UPLOADS_FILE, uploads)
            
            return jsonify({
                "success": True, 
                "message": "教材上傳成功",
                "file_path": f"/static/uploads/{secure_filename}"
            })
        except Exception as e:
            return jsonify({"success": False, "message": f"檔案保存失敗: {str(e)}"})
    else:
        # 處理影片連結
        video_url = request.form.get('videoUrl')
        if not video_url:
            return jsonify({"success": False, "message": "未提供影片連結"})
        
        # 獲取當前日期
        from datetime import datetime
        current_date = datetime.now().strftime("%Y/%m/%d")
        
        # 在真實應用中，這裡會更新資料庫
        new_id = max([p["id"] for p in presentations]) + 1 if presentations else 1
        new_material = {
            "id": new_id,
            "title": title,
            "filename": "video_link",
            "pages": 1,
            "date": current_date,
            "type": "video",
            "video_url": video_url
        }
        presentations.insert(0, new_material)
        save_data_to_file(PRESENTATIONS_FILE, presentations)
        
        return jsonify({
            "success": True, 
            "message": "影片連結新增成功"
        })

# API 路由
@app.route('/api/generate-text', methods=['POST'])
def generate_text():
    prompt = request.json.get('prompt', '')
    
    # 根據提示詞生成回應
    responses = {
        "介紹": "人工智能（AI）是一門讓機器模擬人類智能的科學與技術。它涵蓋機器學習、深度學習、自然語言處理等多個領域。近年來，AI技術快速發展，已經應用於醫療診斷、自動駕駛、智能助手等多個領域，並持續改變著我們的生活和工作方式。",
        "例子": "AI的實際應用例子：\n1. 醫療保健：輔助醫生診斷疾病，分析醫學圖像\n2. 金融服務：欺詐檢測，算法交易，風險評估\n3. 客戶服務：智能聊天機器人，個性化推薦\n4. 教育：自適應學習系統，智能評分\n5. 交通：自動駕駛車輛，交通流量優化",
        "未來": "AI的未來發展趨勢包括：\n- 更強大的通用人工智能（AGI）\n- 人機協作的深入整合\n- 更加透明和可解釋的AI系統\n- 強化隱私和倫理考量的AI設計\n- AI在氣候變化、醫療突破等全球挑戰中的應用"
    }
    
    response = ""
    for key, value in responses.items():
        if key in prompt:
            response = value
            break
    
    if not response:
        response = f"基於您的提示詞「{prompt}」，AI可以生成多種內容。在實際應用中，這裡會連接到GPT等大型語言模型，提供更加精確和相關的回應。這只是一個示例文本，展示AI文字生成功能的界面交互。"
    
    return jsonify({"text": response})

@app.route('/api/generate-image', methods=['POST'])
def generate_image():
    prompt = request.json.get('prompt', '')
    
    # 模擬圖像URL
    categories = ["nature", "city", "tech", "abstract", "people"]
    category = random.choice(categories)
    image_url = f"https://picsum.photos/800/600?category={category}&random={random.randint(1,1000)}"
    
    return jsonify({"image_url": image_url})

@app.route('/api/submit-question', methods=['POST'])
def submit_question():
    try:
        data = request.json
        name = data.get('name', '')
        title = data.get('title', '')
        content = data.get('content', '')
        related_to = data.get('related_to')
        date = data.get('date', datetime.now().strftime('%Y/%m/%d %H:%M'))
        
        # 驗證必要欄位
        if not name or not title or not content:
            return jsonify({"success": False, "message": "缺少必要欄位"})
        
        # 加載現有問題
        global questions
        questions = load_data_from_file(QUESTIONS_FILE, default_questions)
        
        # 生成新問題ID
        new_id = 1
        if questions:
            new_id = max(q['id'] for q in questions) + 1
        
        # 創建新問題
        new_question = {
            "id": new_id,
            "name": name,
            "title": title,
            "content": content,
            "date": date,
            "answer": None,
            "answer_date": None,
            "helpful_count": 0,
            "related_to": related_to,
            "comments": []
        }
        
        # 添加到問題列表
        questions.insert(0, new_question)
        
        # 保存到文件
        success = save_data_to_file(QUESTIONS_FILE, questions)
        if not success:
            return jsonify({"success": False, "message": "保存問題時發生錯誤"})
        
        return jsonify({"success": True, "question": new_question})
    except Exception as e:
        app.logger.error(f"提交問題時發生錯誤: {str(e)}")
        return jsonify({"success": False, "message": f"提交問題時發生錯誤: {str(e)}"})

@app.route('/api/submit-comment', methods=['POST'])
def submit_comment():
    try:
        data = request.json
        question_id = data.get('question_id')
        name = data.get('name', '')
        content = data.get('content', '')
        date = data.get('date', datetime.now().strftime('%Y/%m/%d %H:%M'))
        is_lecturer = session.get('is_lecturer', False)
        
        # 驗證必要欄位
        if not question_id or not name or not content:
            return jsonify({"success": False, "message": "缺少必要欄位"})
        
        # 加載現有問題
        global questions
        questions = load_data_from_file(QUESTIONS_FILE, default_questions)
        
        # 查找問題並添加評論
        for question in questions:
            if question['id'] == int(question_id):
                # 生成評論ID
                comment_id = 1
                if question['comments']:
                    comment_id = max(c['id'] for c in question['comments']) + 1
                
                # 創建新評論
                new_comment = {
                    "id": comment_id,
                    "name": name,
                    "content": content,
                    "date": date,
                    "is_lecturer": is_lecturer
                }
                
                # 添加到評論列表
                question['comments'].append(new_comment)
                
                # 保存到文件
                success = save_data_to_file(QUESTIONS_FILE, questions)
                if not success:
                    return jsonify({"success": False, "message": "保存評論時發生錯誤"})
                
                return jsonify({"success": True, "comment": new_comment})
        
        return jsonify({"success": False, "message": "找不到指定的問題"})
    except Exception as e:
        app.logger.error(f"提交評論時發生錯誤: {str(e)}")
        return jsonify({"success": False, "message": f"提交評論時發生錯誤: {str(e)}"})

@app.route('/api/mark-helpful', methods=['POST'])
def mark_helpful():
    try:
        data = request.json
        question_id = data.get('question_id')
        
        # 驗證必要欄位
        if not question_id:
            return jsonify({"success": False, "message": "缺少必要欄位"})
        
        # 加載現有問題
        global questions
        questions = load_data_from_file(QUESTIONS_FILE, default_questions)
        
        # 查找問題並增加有幫助計數
        for question in questions:
            if question['id'] == int(question_id):
                # 增加計數
                question['helpful_count'] = question.get('helpful_count', 0) + 1
                
                # 保存到文件
                success = save_data_to_file(QUESTIONS_FILE, questions)
                if not success:
                    return jsonify({"success": False, "message": "保存數據時發生錯誤"})
                
                return jsonify({"success": True, "helpful_count": question['helpful_count']})
        
        return jsonify({"success": False, "message": "找不到指定的問題"})
    except Exception as e:
        app.logger.error(f"標記有幫助時發生錯誤: {str(e)}")
        return jsonify({"success": False, "message": f"標記有幫助時發生錯誤: {str(e)}"})

# 下載簡報
@app.route('/download_presentation/<int:presentation_id>')
def download_presentation(presentation_id):
    presentations = load_data_from_file(PRESENTATIONS_FILE, default_presentations)
    
    # 查找對應ID的簡報
    presentation = None
    for p in presentations:
        if p['id'] == presentation_id:
            presentation = p
            break
    
    # 如果找不到簡報，返回404錯誤
    if not presentation or not presentation.get('filename'):
        return render_template('error.html', message="找不到簡報文件"), 404
    
    # 獲取文件路徑
    file_path = os.path.join(app.static_folder, 'uploads', presentation['filename'])
    
    # 如果文件不存在，返回404錯誤
    if not os.path.exists(file_path):
        return render_template('error.html', message="簡報文件不存在"), 404
    
    # 返回文件下載
    return send_from_directory(
        os.path.join(app.static_folder, 'uploads'),
        presentation['filename'],
        as_attachment=True,
        download_name=presentation['filename']
    )

# 刪除教材
@app.route('/lecturer/delete-material', methods=['POST'])
@login_required
def delete_material():
    data = request.get_json()
    material_id = data.get('id')
    material_type = data.get('type', 'presentation')
    
    if not material_id:
        return jsonify({"success": False, "message": "未提供教材ID"})
    
    try:
        material_id = int(material_id)
        
        # 根據類型選擇要操作的列表
        target_list = presentations if material_type == 'presentation' else uploads
        
        # 尋找要刪除的教材
        material_index = None
        material = None
        for i, item in enumerate(target_list):
            if item['id'] == material_id:
                material_index = i
                material = item
                break
        
        if material_index is None:
            return jsonify({"success": False, "message": "找不到指定的教材"})
        
        # 嘗試刪除實際文件
        if 'filename' in material and material['filename'] != 'video_link':
            file_path = os.path.join(os.path.dirname(__file__), 'static', 'uploads', material['filename'])
            if os.path.exists(file_path):
                os.remove(file_path)
        
        # 從列表中移除
        target_list.pop(material_index)
        
        if material_type == 'presentation':
            save_data_to_file(PRESENTATIONS_FILE, presentations)
        else:
            save_data_to_file(UPLOADS_FILE, uploads)
        
        return jsonify({"success": True, "message": "教材已成功刪除"})
    except Exception as e:
        print(f"刪除教材時發生錯誤: {str(e)}")
        return jsonify({"success": False, "message": f"刪除教材時發生錯誤: {str(e)}"})

# 編輯教材
@app.route('/lecturer/edit-material', methods=['POST'])
@login_required
def edit_material():
    try:
        data = request.json
        material_id = data.get('id')
        material_type = data.get('type')
        new_title = data.get('title')
        
        if not material_id or not material_type or not new_title:
            return jsonify({"success": False, "message": "缺少必要參數"})
        
        if material_type == 'presentation':
            presentations = load_data_from_file(PRESENTATIONS_FILE, default_presentations)
            
            # 查找並更新教材
            for presentation in presentations:
                if str(presentation['id']) == str(material_id):
                    presentation['title'] = new_title
                    break
            else:
                return jsonify({"success": False, "message": "找不到指定的簡報"})
            
            # 保存更新後的數據
            if save_data_to_file(PRESENTATIONS_FILE, presentations):
                return jsonify({"success": True, "message": "簡報標題已更新"})
            else:
                return jsonify({"success": False, "message": "保存數據時發生錯誤"})
        
        elif material_type == 'upload':
            uploads = load_data_from_file(UPLOADS_FILE, default_uploads)
            
            # 查找並更新教材
            for upload in uploads:
                if str(upload['id']) == str(material_id):
                    upload['title'] = new_title
                    break
            else:
                return jsonify({"success": False, "message": "找不到指定的上傳教材"})
            
            # 保存更新後的數據
            if save_data_to_file(UPLOADS_FILE, uploads):
                return jsonify({"success": True, "message": "上傳教材標題已更新"})
            else:
                return jsonify({"success": False, "message": "保存數據時發生錯誤"})
        
        else:
            return jsonify({"success": False, "message": "不支持的教材類型"})
    
    except Exception as e:
        return jsonify({"success": False, "message": f"編輯教材時發生錯誤: {str(e)}"})

# 添加參考文獻
@app.route('/lecturer/add-reference', methods=['POST'])
@login_required
def add_reference():
    try:
        title = request.form.get('title')
        url = request.form.get('url')
        
        if not title or not url:
            return jsonify({"success": False, "message": "缺少必要參數"})
        
        # 加載現有參考文獻
        global references
        references = load_data_from_file(REFERENCES_FILE, default_references)
        
        # 生成新的參考文獻ID
        new_id = 1
        if references:
            new_id = max(ref['id'] for ref in references) + 1
        
        # 創建新的參考文獻
        new_reference = {
            "id": new_id,
            "title": title,
            "url": url,
            "date": datetime.now().strftime('%Y/%m/%d')
        }
        
        # 添加到參考文獻列表
        references.append(new_reference)
        
        # 保存更新後的數據
        if save_data_to_file(REFERENCES_FILE, references):
            return jsonify({"success": True, "message": "參考文獻已添加"})
        else:
            return jsonify({"success": False, "message": "保存數據時發生錯誤"})
    
    except Exception as e:
        return jsonify({"success": False, "message": f"添加參考文獻時發生錯誤: {str(e)}"})

# 編輯參考文獻
@app.route('/lecturer/edit-reference', methods=['POST'])
@login_required
def edit_reference():
    try:
        print("收到編輯參考文獻請求")
        data = request.json
        print(f"請求數據: {data}")
        reference_id = data.get('id')
        new_title = data.get('title')
        new_url = data.get('url')
        
        if not reference_id or not new_title or not new_url:
            print(f"缺少必要參數: id={reference_id}, title={new_title}, url={new_url}")
            return jsonify({"success": False, "message": "缺少必要參數"})
        
        # 加載現有參考文獻
        global references
        references = load_data_from_file(REFERENCES_FILE, default_references)
        print(f"當前參考文獻: {references}")
        
        # 查找並更新參考文獻
        for reference in references:
            if str(reference['id']) == str(reference_id):
                print(f"找到參考文獻: {reference}")
                reference['title'] = new_title
                reference['url'] = new_url
                reference['date'] = datetime.now().strftime('%Y/%m/%d')
                print(f"更新後的參考文獻: {reference}")
                break
        else:
            print(f"找不到指定的參考文獻 ID: {reference_id}")
            return jsonify({"success": False, "message": "找不到指定的參考文獻"})
        
        # 保存更新後的數據
        success = save_data_to_file(REFERENCES_FILE, references)
        print(f"保存結果: {success}")
        if success:
            return jsonify({"success": True, "message": "參考文獻已更新"})
        else:
            return jsonify({"success": False, "message": "保存數據時發生錯誤"})
    
    except Exception as e:
        app.logger.error(f"編輯參考文獻時發生錯誤: {str(e)}")
        print(f"編輯參考文獻時發生錯誤: {str(e)}")
        return jsonify({"success": False, "message": f"編輯參考文獻時發生錯誤: {str(e)}"})

# 刪除參考文獻
@app.route('/lecturer/delete-reference', methods=['POST'])
@login_required
def delete_reference():
    try:
        print("收到刪除參考文獻請求")
        data = request.json
        print(f"請求數據: {data}")
        reference_id = data.get('id')
        
        if not reference_id:
            print(f"缺少必要參數: id={reference_id}")
            return jsonify({"success": False, "message": "缺少必要參數"})
        
        # 加載現有參考文獻
        global references
        references = load_data_from_file(REFERENCES_FILE, default_references)
        print(f"當前參考文獻: {references}")
        
        # 查找並刪除參考文獻
        for i, reference in enumerate(references):
            if str(reference['id']) == str(reference_id):
                print(f"找到參考文獻: {reference}")
                del references[i]
                print(f"刪除後的參考文獻列表: {references}")
                break
        else:
            print(f"找不到指定的參考文獻 ID: {reference_id}")
            return jsonify({"success": False, "message": "找不到指定的參考文獻"})
        
        # 保存更新後的數據
        success = save_data_to_file(REFERENCES_FILE, references)
        print(f"保存結果: {success}")
        if success:
            return jsonify({"success": True, "message": "參考文獻已刪除"})
        else:
            return jsonify({"success": False, "message": "保存數據時發生錯誤"})
    
    except Exception as e:
        app.logger.error(f"刪除參考文獻時發生錯誤: {str(e)}")
        print(f"刪除參考文獻時發生錯誤: {str(e)}")
        return jsonify({"success": False, "message": f"刪除參考文獻時發生錯誤: {str(e)}"})

# 添加補充資料
@app.route('/lecturer/add-supplementary', methods=['POST'])
@login_required
def add_supplementary():
    try:
        title = request.form.get('title')
        description = request.form.get('description')
        file = request.files.get('file')
        
        if not title or not description or not file:
            return jsonify({"success": False, "message": "缺少必要參數"})
        
        # 確保文件名安全
        filename = secure_filename(file.filename)
        
        # 檢查文件類型
        if not allowed_file(filename):
            return jsonify({"success": False, "message": "不支持的文件類型"})
        
        # 加載現有補充資料
        global supplementary
        supplementary = load_data_from_file(SUPPLEMENTARY_FILE, default_supplementary)
        
        # 生成新的補充資料ID
        new_id = 1
        if supplementary:
            new_id = max(material['id'] for material in supplementary) + 1
        
        # 保存文件
        file_path = os.path.join(UPLOAD_FOLDER, f"supplementary_{new_id}_{filename}")
        file.save(file_path)
        
        # 創建新的補充資料
        new_supplementary = {
            "id": new_id,
            "title": title,
            "description": description,
            "filename": filename,
            "path": f"uploads/supplementary_{new_id}_{filename}",
            "date": datetime.now().strftime('%Y/%m/%d')
        }
        
        # 添加到補充資料列表
        supplementary.append(new_supplementary)
        
        # 保存更新後的數據
        if save_data_to_file(SUPPLEMENTARY_FILE, supplementary):
            return jsonify({"success": True, "message": "補充資料已添加"})
        else:
            return jsonify({"success": False, "message": "保存數據時發生錯誤"})
    
    except Exception as e:
        return jsonify({"success": False, "message": f"添加補充資料時發生錯誤: {str(e)}"})

# 編輯補充資料
@app.route('/lecturer/edit-supplementary', methods=['POST'])
@login_required
def edit_supplementary():
    try:
        print("收到編輯補充資料請求")
        supplementary_id = request.form.get('id')
        new_title = request.form.get('title')
        new_description = request.form.get('description')
        file = request.files.get('file')
        
        print(f"請求數據: id={supplementary_id}, title={new_title}, description={new_description}, file={file}")
        
        if not supplementary_id or not new_title or not new_description:
            print(f"缺少必要參數: id={supplementary_id}, title={new_title}, description={new_description}")
            return jsonify({"success": False, "message": "缺少必要參數"})
        
        # 加載現有補充資料
        global supplementary
        supplementary = load_data_from_file(SUPPLEMENTARY_FILE, default_supplementary)
        print(f"當前補充資料: {supplementary}")
        
        # 查找並更新補充資料
        for material in supplementary:
            if str(material['id']) == str(supplementary_id):
                print(f"找到補充資料: {material}")
                material['title'] = new_title
                material['description'] = new_description
                material['date'] = datetime.now().strftime('%Y/%m/%d')
                
                # 如果有新文件，則更新文件
                if file and file.filename:
                    print(f"有新文件: {file.filename}")
                    # 確保文件名安全
                    filename = secure_filename(file.filename)
                    
                    # 檢查文件類型
                    if not allowed_file(filename):
                        print(f"不支持的文件類型: {filename}")
                        return jsonify({"success": False, "message": "不支持的文件類型"})
                    
                    # 保存文件
                    file_path = os.path.join(UPLOAD_FOLDER, f"supplementary_{supplementary_id}_{filename}")
                    print(f"保存文件到: {file_path}")
                    file.save(file_path)
                    
                    # 更新文件信息
                    material['filename'] = filename
                    material['path'] = f"uploads/supplementary_{supplementary_id}_{filename}"
                
                print(f"更新後的補充資料: {material}")
                break
        else:
            print(f"找不到指定的補充資料 ID: {supplementary_id}")
            return jsonify({"success": False, "message": "找不到指定的補充資料"})
        
        # 保存更新後的數據
        success = save_data_to_file(SUPPLEMENTARY_FILE, supplementary)
        print(f"保存結果: {success}")
        if success:
            return jsonify({"success": True, "message": "補充資料已更新"})
        else:
            return jsonify({"success": False, "message": "保存數據時發生錯誤"})
    
    except Exception as e:
        app.logger.error(f"編輯補充資料時發生錯誤: {str(e)}")
        print(f"編輯補充資料時發生錯誤: {str(e)}")
        return jsonify({"success": False, "message": f"編輯補充資料時發生錯誤: {str(e)}"})

# 刪除補充資料
@app.route('/lecturer/delete-supplementary', methods=['POST'])
@login_required
def delete_supplementary():
    try:
        print("收到刪除補充資料請求")
        data = request.json
        print(f"請求數據: {data}")
        supplementary_id = data.get('id')
        
        if not supplementary_id:
            print(f"缺少必要參數: id={supplementary_id}")
            return jsonify({"success": False, "message": "缺少必要參數"})
        
        # 加載現有補充資料
        global supplementary
        supplementary = load_data_from_file(SUPPLEMENTARY_FILE, default_supplementary)
        print(f"當前補充資料: {supplementary}")
        
        # 查找並刪除補充資料
        for i, material in enumerate(supplementary):
            if str(material['id']) == str(supplementary_id):
                print(f"找到補充資料: {material}")
                # 刪除文件
                if material.get('path'):
                    file_path = os.path.join(app.static_folder, material['path'])
                    if os.path.exists(file_path):
                        print(f"刪除文件: {file_path}")
                        os.remove(file_path)
                
                del supplementary[i]
                print(f"刪除後的補充資料列表: {supplementary}")
                break
        else:
            print(f"找不到指定的補充資料 ID: {supplementary_id}")
            return jsonify({"success": False, "message": "找不到指定的補充資料"})
        
        # 保存更新後的數據
        success = save_data_to_file(SUPPLEMENTARY_FILE, supplementary)
        print(f"保存結果: {success}")
        if success:
            return jsonify({"success": True, "message": "補充資料已刪除"})
        else:
            return jsonify({"success": False, "message": "保存數據時發生錯誤"})
    
    except Exception as e:
        return jsonify({"success": False, "message": f"刪除補充資料時發生錯誤: {str(e)}"})

@app.route('/get_slide_image/<presentation_id>/<page_number>')
def get_slide_image(presentation_id, page_number):
    """獲取特定簡報的特定頁面圖片，用於全螢幕模式下的AJAX請求"""
    try:
        # 獲取簡報信息
        presentation = get_presentation_by_id(presentation_id)
        if not presentation:
            return jsonify({'success': False, 'error': '找不到簡報'})
        
        # 檢查頁碼是否有效
        page_number = int(page_number)
        if page_number < 1 or page_number > presentation.total_pages:
            return jsonify({'success': False, 'error': '頁碼超出範圍'})
        
        # 構建圖片URL
        image_url = f'/static/slides/presentation-{presentation_id}/slide-{page_number}.png'
        
        return jsonify({
            'success': True, 
            'image_url': image_url,
            'current_page': page_number,
            'total_pages': presentation.total_pages
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/api/get-question-comments', methods=['GET'])
def get_question_comments():
    try:
        question_id = request.args.get('question_id')
        
        if not question_id:
            return jsonify({"success": False, "message": "缺少問題ID"})
        
        # 加載現有問題
        global questions
        questions = load_data_from_file(QUESTIONS_FILE, default_questions)
        
        # 查找問題並獲取評論
        for question in questions:
            if question['id'] == int(question_id):
                return jsonify({"success": True, "comments": question.get('comments', [])})
        
        return jsonify({"success": False, "message": "找不到指定的問題"})
    except Exception as e:
        app.logger.error(f"獲取問題評論時發生錯誤: {str(e)}")
        return jsonify({"success": False, "message": f"獲取問題評論時發生錯誤: {str(e)}"})

@app.before_request
def load_lecturer_info():
    # 跳過靜態文件請求
    if request.path.startswith('/static/'):
        return
    
    # 只在未設置講師資料但已經登入的情況下載入持久化資料
    if 'lecturer_username' in session and ('lecturer_name' not in session or 'lecturer_title' not in session or 'lecturer_email' not in session):
        # 嘗試從配置文件讀取講師資料
        config_path = os.path.join(os.path.dirname(__file__), 'static', 'config', 'lecturer_info.json')
        
        if os.path.exists(config_path):
            try:
                with open(config_path, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    
                    # 如果配置文件中有講師資料，則載入到session
                    if 'lecturer_name' in config:
                        session['lecturer_name'] = config.get('lecturer_name')
                    if 'lecturer_title' in config:
                        session['lecturer_title'] = config.get('lecturer_title')
                    if 'lecturer_email' in config:
                        session['lecturer_email'] = config.get('lecturer_email')
                        print("已從配置文件載入講師資料")
            except Exception as e:
                print(f"載入講師資料時發生錯誤: {str(e)}")

if __name__ == '__main__':
    # 確保資料夾存在
    os.makedirs('static/downloads', exist_ok=True)
    app.run(debug=True, host='0.0.0.0', port=8080)